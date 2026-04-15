from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DESKTOP = Path.home() / "Desktop"
OUTPUT = DESKTOP / "室内视觉引导系统周汇报_20260327.pptx"

BG = RGBColor(245, 247, 250)
NAVY = RGBColor(18, 42, 66)
TEAL = RGBColor(23, 126, 137)
ORANGE = RGBColor(238, 155, 0)
RED = RGBColor(196, 69, 54)
GREEN = RGBColor(65, 145, 108)
GRAY = RGBColor(90, 102, 117)
LIGHT = RGBColor(255, 255, 255)
LINE = RGBColor(214, 220, 227)
ACCENT = RGBColor(228, 236, 242)


def set_font(run, size, color, bold=False, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.6), Inches(-0.2), Inches(4.1), Inches(1.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.fill.transparency = 0.15
    shape.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.8), Inches(5.8), Inches(2.4), Inches(2.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(225, 239, 242)
    circle.fill.transparency = 0.2
    circle.line.fill.background()


def add_footer(slide, text="代码依据：Core/Src/main.c  Core/Src/oled.c"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, 9, GRAY)
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, 26, NAVY, bold=True)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.25), Inches(0.5), Inches(1.45), Inches(0.42)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = TEAL
    tag.line.fill.background()
    tf2 = tag.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "周进展"
    set_font(r2, 12, LIGHT, bold=True)
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.18), Inches(8.5), Inches(0.4))
        tf3 = sub.text_frame
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = subtitle
        set_font(r3, 12, GRAY)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.58), Inches(1.7), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()


def add_bullet_block(slide, left, top, width, height, items, title=None, fill_rgb=LIGHT):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_rgb
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    tx = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.18), width - Inches(0.44), height - Inches(0.25))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if title:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_font(r, 16, NAVY, bold=True)
        for item in items:
            p = tf.add_paragraph()
            p.level = 0
            p.space_before = Pt(4)
            r = p.add_run()
            r.text = f"• {item}"
            set_font(r, 13, GRAY)
    else:
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            p.space_before = Pt(4)
            r = p.add_run()
            r.text = f"• {item}"
            set_font(r, 13, GRAY)


def add_component_box(slide, left, top, width, height, title, desc, color):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    set_font(r1, 17, color, bold=True)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = desc
    set_font(r2, 11, GRAY)


def add_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()


def add_formula_box(slide, left, top, width, height, title, formula, note):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_font(r1, 15, NAVY, bold=True)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = formula
    set_font(r2, 18, TEAL, bold=True, name="Consolas")

    p3 = tf.add_paragraph()
    p3.space_before = Pt(8)
    r3 = p3.add_run()
    r3.text = note
    set_font(r3, 12, GRAY)


def add_status_card(slide, left, top, width, height, title, lines, color):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = color
    card.line.width = Pt(2)

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + Inches(0.18), top + Inches(0.16), Inches(0.95), Inches(0.32)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf_badge = badge.text_frame
    tf_badge.clear()
    p_badge = tf_badge.paragraphs[0]
    p_badge.alignment = PP_ALIGN.CENTER
    r_badge = p_badge.add_run()
    r_badge.text = "已完成" if color == GREEN else "待优化"
    set_font(r_badge, 10, LIGHT, bold=True)
    tf_badge.vertical_anchor = MSO_ANCHOR.MIDDLE

    tx = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.58), width - Inches(0.4), height - Inches(0.7))
    tf = tx.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_font(r1, 16, NAVY, bold=True)
    for line in lines:
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"• {line}"
        set_font(r, 12, GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.85), Inches(7.2), Inches(4.4))
    hero.fill.solid()
    hero.fill.fore_color.rgb = LIGHT
    hero.line.fill.background()

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.1), Inches(2.5), Inches(0.42))
    tag.fill.solid()
    tag.fill.fore_color.rgb = TEAL
    tag.line.fill.background()
    tf_tag = tag.text_frame
    tf_tag.clear()
    p_tag = tf_tag.paragraphs[0]
    p_tag.alignment = PP_ALIGN.CENTER
    r_tag = p_tag.add_run()
    r_tag.text = "OpenMV + STM32 项目"
    set_font(r_tag, 12, LIGHT, bold=True)
    tf_tag.vertical_anchor = MSO_ANCHOR.MIDDLE

    title = slide.shapes.add_textbox(Inches(0.95), Inches(1.75), Inches(6.5), Inches(1.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "室内目标位置引导系统\n周进展汇报"
    set_font(r, 26, NAVY, bold=True)

    sub = slide.shapes.add_textbox(Inches(0.98), Inches(3.5), Inches(5.5), Inches(0.8))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "围绕图像坐标传输、坐标解析、OLED 调试显示\n以及舵机左右跟随完成阶段性联调"
    set_font(r2, 14, GRAY)

    date = slide.shapes.add_textbox(Inches(1.0), Inches(4.55), Inches(3.0), Inches(0.4))
    tf3 = date.text_frame
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "汇报日期：2026.03.27"
    set_font(r3, 12, ORANGE, bold=True)

    for i, (name, color, x, y) in enumerate(
        [
            ("OpenMV", TEAL, 8.55, 1.2),
            ("UART", ORANGE, 10.1, 1.9),
            ("STM32", NAVY, 8.8, 3.1),
            ("OLED", GREEN, 10.2, 4.0),
            ("SERVO", RED, 8.6, 4.85),
        ]
    ):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(0.62))
        shp.fill.solid()
        shp.fill.fore_color.rgb = LIGHT
        shp.line.color.rgb = color
        shp.line.width = Pt(2)
        tfb = shp.text_frame
        tfb.clear()
        pb = tfb.paragraphs[0]
        pb.alignment = PP_ALIGN.CENTER
        rb = pb.add_run()
        rb.text = name
        set_font(rb, 14, color, bold=True)
        tfb.vertical_anchor = MSO_ANCHOR.MIDDLE

    link1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(9.0), Inches(2.15), Inches(0.9), Inches(0.35))
    link1.fill.solid()
    link1.fill.fore_color.rgb = ACCENT
    link1.line.fill.background()
    link2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(9.25), Inches(4.28), Inches(0.9), Inches(0.35))
    link2.fill.solid()
    link2.fill.fore_color.rgb = ACCENT
    link2.line.fill.background()
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "课题背景与项目目标", "解决“到附近了，但还找不到具体房间”的近距离引导问题")
    add_bullet_block(
        slide,
        Inches(0.7),
        Inches(1.95),
        Inches(5.95),
        Inches(3.9),
        [
            "室内环境中，用户即使已经到达目标区域附近，仍然可能难以快速锁定具体房间或门牌。",
            "相比纯文字或平面地图提示，视觉与方向结合的引导方式更直观，更适合近距离定位。",
            "本课题希望把“识别目标位置 + 判断左右方向 + 执行基础引导”做成可落地的嵌入式系统。",
        ],
        title="课题背景",
    )
    add_bullet_block(
        slide,
        Inches(6.95),
        Inches(1.95),
        Inches(5.7),
        Inches(3.9),
        [
            "搭建 OpenMV、STM32、OLED、舵机之间的完整数据链路。",
            "让系统可以接收目标坐标、解析位置，并驱动云台进行左右方向跟随。",
            "为后续门牌识别、路径提示、语音交互等功能扩展打下基础。",
        ],
        title="项目目标",
        fill_rgb=RGBColor(251, 252, 254),
    )
    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(6.15), Inches(11.9), Inches(0.7))
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(232, 244, 246)
    note.line.fill.background()
    tf = note.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "本周工作重点：先把坐标通信、显示调试和左右跟随跑通，完成系统基础闭环。"
    set_font(r, 14, NAVY, bold=True)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "系统组成与整体链路", "目前已经具备“采集 - 传输 - 处理 - 执行 - 显示”这条主链路")
    add_component_box(slide, Inches(0.7), Inches(2.2), Inches(2.1), Inches(1.2), "OpenMV", "采集图像\n输出目标坐标", TEAL)
    add_arrow(slide, Inches(2.95), Inches(2.55), Inches(0.7), Inches(0.45))
    add_component_box(slide, Inches(3.7), Inches(2.2), Inches(1.9), Inches(1.2), "UART", "串口通信\n发送/接收数据", ORANGE)
    add_arrow(slide, Inches(5.75), Inches(2.55), Inches(0.7), Inches(0.45))
    add_component_box(slide, Inches(6.45), Inches(2.2), Inches(2.2), Inches(1.2), "STM32", "接收数据\n解析与控制", NAVY)
    add_arrow(slide, Inches(8.8), Inches(2.55), Inches(0.7), Inches(0.45))
    add_component_box(slide, Inches(9.55), Inches(1.45), Inches(2.25), Inches(1.15), "OLED", "显示坐标\n运行状态", GREEN)
    add_component_box(slide, Inches(9.55), Inches(3.0), Inches(2.25), Inches(1.15), "舵机云台", "PWM 控制\n左右引导", RED)

    add_bullet_block(
        slide,
        Inches(0.85),
        Inches(4.35),
        Inches(11.8),
        Inches(1.9),
        [
            "OpenMV 端负责识别目标并给出坐标。",
            "STM32 使用 USART2 中断接收数据，并完成坐标解析、OLED 刷新和舵机控制。",
            "OLED 用于观察 X/Y 调试信息，舵机云台负责输出左右方向指示。",
        ],
        title="模块分工",
    )
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "本周已完成内容", "根据当前代码，已经完成的核心功能如下")
    add_status_card(
        slide,
        Inches(0.8),
        Inches(1.95),
        Inches(3.8),
        Inches(2.0),
        "通信链路",
        [
            "OpenMV 串口发送坐标数据",
            "STM32 串口中断逐字节接收",
            "接收到一帧后置位标志进行处理",
        ],
        GREEN,
    )
    add_status_card(
        slide,
        Inches(4.8),
        Inches(1.95),
        Inches(3.8),
        Inches(2.0),
        "数据处理",
        [
            "使用 sscanf 解析 X/Y 坐标",
            "判断新旧坐标是否变化",
            "实现基础调试数据缓存与更新",
        ],
        GREEN,
    )
    add_status_card(
        slide,
        Inches(8.8),
        Inches(1.95),
        Inches(3.8),
        Inches(2.0),
        "输出执行",
        [
            "OLED 实时显示 X 与 Y",
            "TIM2 PWM 驱动舵机",
            "云台可根据目标位置左右跟随",
        ],
        GREEN,
    )
    add_bullet_block(
        slide,
        Inches(1.0),
        Inches(4.45),
        Inches(11.5),
        Inches(1.55),
        [
            "本周最重要的结果不是单个模块，而是系统整体链路已经从“识别坐标”打通到了“舵机动作 + OLED 显示”。"
        ],
        title="阶段性结论",
        fill_rgb=RGBColor(251, 252, 254),
    )
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "关键代码实现说明", "内容直接对应当前 main.c 与 oled.c 的真实逻辑")
    add_formula_box(
        slide,
        Inches(0.8),
        Inches(1.95),
        Inches(4.05),
        Inches(1.8),
        "串口协议与解析",
        "X:120,Y:64",
        "USART2 中断接收，遇到回车换行认为一帧完成，再使用 sscanf 解析坐标。",
    )
    add_formula_box(
        slide,
        Inches(0.8),
        Inches(4.0),
        Inches(4.05),
        Inches(1.8),
        "舵机角度映射",
        "target = (x / 160.0) * 180.0",
        "把图像中的 X 坐标映射到舵机转角范围，实现左右方向跟随。",
    )
    add_formula_box(
        slide,
        Inches(4.98),
        Inches(1.95),
        Inches(4.05),
        Inches(1.8),
        "当前平滑控制",
        "angle = 0.8 * old + 0.2 * target",
        "使用一次简单低通方式减弱跳变，再把角度限制在 10° 到 170°。",
    )
    add_formula_box(
        slide,
        Inches(4.98),
        Inches(4.0),
        Inches(4.05),
        Inches(1.8),
        "PWM 输出参数",
        "500us ~ 2500us / 50Hz",
        "TIM2 输出 PWM，占空控制对应舵机角度，已经可以稳定动作。",
    )
    add_bullet_block(
        slide,
        Inches(9.15),
        Inches(1.95),
        Inches(3.45),
        Inches(3.85),
        [
            "OLED 初始化成功",
            "显示 X / Y 调试值",
            "坐标变化时刷新屏幕",
            "有利于联调与排错",
        ],
        title="OLED 调试作用",
        fill_rgb=RGBColor(251, 252, 254),
    )
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "当前效果与存在问题", "系统已经能跑起来，但控制品质还需要继续打磨")
    add_status_card(
        slide,
        Inches(0.85),
        Inches(1.95),
        Inches(3.85),
        Inches(2.1),
        "当前效果",
        [
            "坐标可以实时传输和显示",
            "舵机能够根据目标位置左右跟随",
            "OpenMV 到执行端的整体链路已打通",
        ],
        GREEN,
    )
    add_status_card(
        slide,
        Inches(4.9),
        Inches(1.95),
        Inches(3.85),
        Inches(2.1),
        "当前问题",
        [
            "舵机跟随时仍然存在抖动",
            "控制方式偏简单，抗干扰不足",
            "连续运行时的系统稳定性还需验证",
        ],
        RED,
    )
    add_status_card(
        slide,
        Inches(8.95),
        Inches(1.95),
        Inches(3.4),
        Inches(2.1),
        "原因判断",
        [
            "目标坐标波动直接影响舵机目标角度",
            "当前只做了基础平滑，尚未加入死区",
            "尚未做更系统的滤波和控制优化",
        ],
        ORANGE,
    )
    add_bullet_block(
        slide,
        Inches(1.0),
        Inches(4.55),
        Inches(11.4),
        Inches(1.45),
        [
            "阶段评价：功能上已经实现“能用”，下一步重点是把“能用”进一步优化成“更稳、更顺”。"
        ],
        title="本周结论",
    )
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "后续计划", "下一阶段将围绕跟随效果、控制品质和功能扩展继续推进")
    add_bullet_block(
        slide,
        Inches(0.85),
        Inches(2.0),
        Inches(5.6),
        Inches(3.7),
        [
            "优化 X 轴跟随效果，减少云台抖动。",
            "加入死区和平滑控制，避免小范围波动频繁触发动作。",
            "继续提升串口接收与整体运行稳定性。",
            "完善调试信息，便于后续联调和问题定位。",
        ],
        title="短期计划",
    )
    add_bullet_block(
        slide,
        Inches(6.75),
        Inches(2.0),
        Inches(5.6),
        Inches(3.7),
        [
            "扩展门牌识别，提高目标定位能力。",
            "尝试加入语音交互，增强用户提示体验。",
            "结合更完善的控制策略，提升系统实用性。",
            "逐步向完整的室内近距离引导原型推进。",
        ],
        title="中期方向",
        fill_rgb=RGBColor(251, 252, 254),
    )
    end = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(6.15), Inches(11.0), Inches(0.65))
    end.fill.solid()
    end.fill.fore_color.rgb = RGBColor(232, 244, 246)
    end.line.fill.background()
    tf = end.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "汇报总结：本周已经完成基础闭环搭建，接下来进入“控制优化 + 功能扩展”阶段。"
    set_font(r, 14, NAVY, bold=True)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_footer(slide)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
